import streamlit as st
import tempfile
import base64
from PIL import Image 
from pathlib import Path
import os 
import logging
import streamlit.components.v1 as components # For st.components.v1.iframe
from user_management import get_user_dir
import datetime
import auth
auth.initialize_session()

# Assuming these functions are correctly defined in utils.py
# and that they return Document objects with appropriate metadata, including "title" where applicable.
from utils import (
    image_to_document, url_to_document, youtube_to_document, extract_video_id,
    render_footer
)
# LlamaIndex imports
from llama_index.core import VectorStoreIndex, Document as LlamaDocument, get_response_synthesizer
# Your RAG Engine
from Tutor_chat import RAGQueryEngine 

def render():
    """Renders the research mode page with enhanced content previews and title display."""
    st.title("🔬 Research Mode")
    st.markdown("Upload documents, images, or provide URLs. Preview your content, then build an index to chat with it.")

    # --- Initialize session state variables ---
    if "research_items_for_preview_and_indexing" not in st.session_state:
        st.session_state.research_items_for_preview_and_indexing = []
    if "research_processed_source_ids" not in st.session_state: 
        st.session_state.research_processed_source_ids = set()
    if "research_index" not in st.session_state:
        st.session_state.research_index = None
    if "research_chat_history" not in st.session_state:
        st.session_state.research_chat_history = []

    # --- UI for Uploading Content ---
    with st.expander("📁 Upload Files (PDF, DOCX, PPTX, TXT, Images)", expanded=True):
        uploaded_files = st.file_uploader(
            "Choose files (PDFs, text files, images, etc.)", 
            type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg"],
            accept_multiple_files=True, key="research_file_uploader_final"
        )
        if uploaded_files:
            for uploaded_file in uploaded_files:
                file_id = ("file", uploaded_file.name, uploaded_file.size) 

                if file_id not in st.session_state.research_processed_source_ids:
                    tmp_file_path = None
                    try:
                        # Save to a temporary file to get a stable path for processing
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{uploaded_file.name}") as tmp_file:
                            tmp_file.write(uploaded_file.getvalue())
                            tmp_file_path = tmp_file.name
                        
                        text_content = ""
                        preview_type = "text" # Default
                        preview_content = None # Will hold bytes or text for preview
                        display_source_name = uploaded_file.name # Default display name
                        
                        file_ext = Path(uploaded_file.name).suffix.lower()

                        if file_ext == ".pdf":
                            import fitz # PyMuPDF
                            with fitz.open(tmp_file_path) as doc_pdf:
                                text_content = "".join(page.get_text() for page in doc_pdf)
                            preview_type = "pdf"
                            preview_content = uploaded_file.getvalue() 
                        elif file_ext in [".png", ".jpg", ".jpeg"]:
                            # For images, OCR for text, but keep original image for preview
                            # image_to_document should ideally take bytes or path.
                            # For consistency, we pass the BytesIO object (uploaded_file)
                            doc_from_img = image_to_document(uploaded_file) 
                            text_content = doc_from_img.text
                            preview_type = "image"
                            preview_content = uploaded_file.getvalue()
                        elif file_ext == ".txt":
                            text_content = Path(tmp_file_path).read_text(encoding="utf-8", errors="replace")
                            preview_type = "text"
                            preview_content = text_content
                        elif file_ext == ".docx":
                            from docx import Document as DocxFile # Local import for clarity
                            doc_obj = DocxFile(tmp_file_path)
                            text_content = "\n".join([para.text for para in doc_obj.paragraphs])
                            preview_type = "text"
                            preview_content = text_content
                        elif file_ext == ".pptx":
                            from pptx import Presentation # Local import
                            pres = Presentation(tmp_file_path)
                            slide_texts = []
                            for slide in pres.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        slide_texts.append(shape.text)
                            text_content = "\n".join(slide_texts)
                            preview_type = "text"
                            preview_content = text_content
                        
                        if text_content.strip() or preview_type in ["image", "pdf"]:
                            doc_metadata = {"source": uploaded_file.name, "size": uploaded_file.size, "id_tuple": str(file_id), "title": uploaded_file.name}
                            item_doc = LlamaDocument(text=text_content if text_content.strip() else "No text extracted (visual preview only)", metadata=doc_metadata)
                            
                            st.session_state.research_items_for_preview_and_indexing.append({
                                "doc": item_doc,
                                "display_source_name": display_source_name,
                                "preview_type": preview_type,
                                "preview_content": preview_content
                            })
                            st.session_state.research_processed_source_ids.add(file_id)
                            st.success(f"Processed '{uploaded_file.name}' for research context.")
                        else:
                            st.warning(f"Could not extract meaningful content or not a supported preview type for '{uploaded_file.name}'.")
                    
                    except Exception as e:
                        st.error(f"Error processing file {uploaded_file.name}: {e}")
                        logging.error(f"File processing error for {uploaded_file.name}: {e}", exc_info=True)
                    finally:
                        if tmp_file_path and os.path.exists(tmp_file_path):
                            os.unlink(tmp_file_path) # Ensure temp file cleanup
                else:
                    st.info(f"'{uploaded_file.name}' has already been processed in this session.")

    # Web Page URL
    with st.expander("🔗 Add Web Page URL", expanded=False):
        web_url_input = st.text_input("Enter URL:", key="research_web_url_final")
        if st.button("Fetch Web Content", key="fetch_web_final_button"):
            if web_url_input:
                url_id = ("url", web_url_input)
                if url_id not in st.session_state.research_processed_source_ids:
                    with st.spinner(f"Fetching content from {web_url_input}..."):
                        doc_from_url = url_to_document(web_url_input) 
                        if not doc_from_url.metadata.get("error"):
                            doc_metadata = doc_from_url.metadata.copy() # metadata includes 'title' from utils.py
                            doc_metadata.update({"id_tuple": str(url_id)})
                            doc_from_url.metadata = doc_metadata
                            
                            display_name = doc_from_url.metadata.get("title", web_url_input)

                            st.session_state.research_items_for_preview_and_indexing.append({
                                "doc": doc_from_url,
                                "display_source_name": display_name,
                                "preview_type": "webpage",
                                "preview_content": web_url_input 
                            })
                            st.session_state.research_processed_source_ids.add(url_id)
                            st.success(f"Added web content: '{display_name}'.")
                        else:
                            st.error(f"Failed to fetch URL: {doc_from_url.text}")
                else:
                    st.info(f"URL '{web_url_input}' already processed.")
            else:
                st.warning("Please enter a URL.")

    # YouTube URL
    with st.expander("📺 Add YouTube Video URL for Transcript", expanded=False):
        yt_url_input_str = st.text_input("Enter YouTube video URL:", key="research_yt_url_final")
        if st.button("Fetch Transcript & Add Video", key="fetch_yt_final_button"):
            if yt_url_input_str:
                video_id = extract_video_id(yt_url_input_str)
                if not video_id:
                    st.error("Invalid YouTube URL. Could not extract Video ID.")
                else:
                    yt_id_tuple = ("youtube", video_id)
                    if yt_id_tuple not in st.session_state.research_processed_source_ids:
                        with st.spinner(f"Fetching transcript for {yt_url_input_str}..."):
                            doc_from_yt = youtube_to_document(yt_url_input_str) 
                            if not doc_from_yt.metadata.get("error"):
                                doc_metadata = doc_from_yt.metadata.copy() # metadata includes 'title' from utils.py
                                doc_metadata.update({"id_tuple": str(yt_id_tuple)})
                                doc_from_yt.metadata = doc_metadata
                                
                                display_name = doc_from_yt.metadata.get("title", f"YouTube Video: {video_id}")
                                # Use the original input URL or a constructed watch URL for embedding
                                embed_url = doc_from_yt.metadata.get("source", f"https://www.youtube.com/watch?v={video_id}")

                                st.session_state.research_items_for_preview_and_indexing.append({
                                    "doc": doc_from_yt, 
                                    "display_source_name": display_name,
                                    "preview_type": "youtube",
                                    "preview_content": embed_url 
                                })
                                st.session_state.research_processed_source_ids.add(yt_id_tuple)
                                st.success(f"Added transcript: '{display_name}'.")
                            else:
                                st.error(f"Failed to fetch transcript: {doc_from_yt.text}")
                    else:
                        st.info(f"YouTube video for '{video_id}' already processed.")
            else:
                st.warning("Please enter a YouTube URL.")
    
    st.markdown("---")

    # --- Display Processed Documents & Previews ---
    if st.session_state.research_items_for_preview_and_indexing:
        st.subheader(f"📚 Content Added for Research ({len(st.session_state.research_items_for_preview_and_indexing)} item(s))")
        
        for i, item_data in enumerate(st.session_state.research_items_for_preview_and_indexing):
            doc_object = item_data["doc"]
            display_name = item_data["display_source_name"] 
            preview_type = item_data["preview_type"]
            preview_content = item_data["preview_content"]

            with st.expander(f"{i+1}. {display_name}", expanded=False): # Display name used here
                if preview_type == "image":
                    st.image(preview_content, caption=display_name, use_column_width=True)
                elif preview_type == "pdf":
                    st.info("PDF Preview: For interactive PDF viewing, ensure the browser supports it or consider downloading.")
                    # Display PDF using st.pdf_viewer (Streamlit native) or base64 iframe
                    # st.pdf_viewer requires bytes. preview_content for PDF is already bytes.
                    try:
                        st.pdf_viewer(preview_content, height=500) # Streamlit 1.33+
                    except AttributeError: # Fallback for older Streamlit or if pdf_viewer fails
                        st.warning("`st.pdf_viewer` not available or failed. Trying iframe for PDF.")
                        b64_pdf = base64.b64encode(preview_content).decode('utf-8')
                        components.iframe(f"data:application/pdf;base64,{b64_pdf}#toolbar=0&navpanes=0", height=500, scrolling=True)

                elif preview_type == "webpage":
                    st.markdown(f"**Web Page Preview for:** [{display_name}]({preview_content})")
                    try:
                        components.iframe(preview_content, height=400, scrolling=True)
                        st.caption("Note: Some websites may not allow embedding due to security policies (X-Frame-Options).")
                    except Exception as e_iframe:
                        st.warning(f"Could not embed iframe for {preview_content}: {e_iframe}. Link provided above.")
                elif preview_type == "youtube":
                    st.markdown(f"**YouTube Video Preview for:** [{display_name}]({preview_content})") # Display_name is now better
                    st.video(preview_content) # preview_content is the embeddable URL
                
                # Always show extracted text (if any) for indexing transparency
                doc_text_for_preview = doc_object.text if doc_object else "Document object missing."
                if doc_text_for_preview and doc_text_for_preview != "No text extracted (visual preview only)":
                    st.text_area("Text Content for Indexing (preview)", doc_text_for_preview[:1000]+"...", height=150, disabled=True, key=f"text_preview_final_{i}")
                elif preview_type not in ["image", "pdf", "youtube", "webpage"] and preview_content and isinstance(preview_content, str):
                     st.text_area("Text Content", preview_content[:1000]+"...", height=150, disabled=True, key=f"text_preview_other_final_{i}")
                elif not doc_text_for_preview or doc_text_for_preview == "No text extracted (visual preview only)":
                    st.caption("No significant text content was extracted for indexing for this item (visual content prioritized).")


        if st.button("Prepare Research Context from Added Content", key="prepare_research_context_final_button", use_container_width=True):
            if not st.session_state.research_items_for_preview_and_indexing:
                st.warning("No content has been added to prepare the research context.")
            else:
                with st.spinner("Creating research index... This may take a moment."):
                    docs_for_indexing = [item["doc"] for item in st.session_state.research_items_for_preview_and_indexing]
                    st.session_state.research_index = VectorStoreIndex.from_documents(docs_for_indexing)
                    st.session_state.research_chat_history = [] 
                    st.success("Research context is ready! You can now ask questions below.")
                    st.rerun() 
    
    # --- Chat Interface (Appears if index is built) ---
    if st.session_state.research_index:
        st.subheader("💬 Chat with Your Research Content")
        for msg_entry in st.session_state.research_chat_history:
            role = msg_entry.get("role")
            content = msg_entry.get("content", "")
            with st.chat_message(role): st.markdown(content)
        
        user_query_research = st.chat_input("Ask a question about your research content...", key="research_chat_input_final")
        if user_query_research:
            st.session_state.research_chat_history.append({"role": "user", "content": user_query_research})
            with st.chat_message("user"): st.markdown(user_query_research)
            with st.spinner("Thinking..."):
                research_retriever = st.session_state.research_index.as_retriever(similarity_top_k=3)
                research_synthesizer = get_response_synthesizer(response_mode="compact")
                query_engine_research = RAGQueryEngine(retriever=research_retriever, response_synthesizer=research_synthesizer, mode="research" )
                response_text = query_engine_research.custom_query(user_query_research)
            st.session_state.research_chat_history.append({"role": "assistant", "content": response_text})
            with st.chat_message("assistant"): st.markdown(response_text)

        if st.session_state.research_chat_history:
            export_text = "\n\n".join(f"{entry['role']}: {entry['content']}" for entry in st.session_state.research_chat_history)
            file_name = f"research_chat_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
            user_id = st.session_state.get("user_name")
            if user_id:
                research_dir = os.path.join(get_user_dir(user_id), "research")
                os.makedirs(research_dir, exist_ok=True)
                file_path = os.path.join(research_dir, file_name)
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(export_text)
                st.success(f"Chat log saved to your research folder as {file_name}.")

            st.download_button("📄 Download Chat", data=export_text, file_name=file_name, mime="text/plain", key="download_research_chat_final")

    elif not st.session_state.research_items_for_preview_and_indexing: 
        st.info("Upload or add content using the sections above to start your research session.")
    
    if st.session_state.research_items_for_preview_and_indexing or st.session_state.research_index:
        if st.button("Clear All Research Content & Chat", key="clear_research_all_final_button", type="secondary"):
            st.session_state.research_index = None
            st.session_state.research_items_for_preview_and_indexing = []
            st.session_state.research_chat_history = []
            st.session_state.research_processed_source_ids = set()
            st.success("Research context and chat history cleared.")
            st.rerun()

    render_footer()