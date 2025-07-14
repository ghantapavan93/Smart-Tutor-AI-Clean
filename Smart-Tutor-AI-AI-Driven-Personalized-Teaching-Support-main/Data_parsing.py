import os
import torch
import re
import unicodedata
import pdfplumber
from pathlib import Path
from pptx import Presentation
from transformers import AutoModel, AutoTokenizer
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings, get_response_synthesizer
from llama_index.core.ingestion import IngestionPipeline
from llama_index.core.node_parser import SentenceSplitter, SemanticSplitterNodeParser
try:
    from llama_index.core.node_parser import CodeSplitter
    CODE_SPLITTER_AVAILABLE = True
except ImportError:
    CODE_SPLITTER_AVAILABLE = False
    print("⚠️ CodeSplitter not available, using standard splitting for code")
from llama_index.vector_stores.chroma import ChromaVectorStore
from chromadb import PersistentClient
from llama_index.core.retrievers import BaseRetriever
from llama_index.core.query_engine import CustomQueryEngine
from llama_index.core.response_synthesizers import BaseSynthesizer
from llama_index.llms.ollama import Ollama
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.readers.file import IPYNBReader

# ------------------------------
# CUSTOM PPTX READER
# ------------------------------
class PPTXTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        prs = Presentation(file_path)
        docs: list[Document] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw = shape.text.strip()
                    if raw:
                        text_runs.append(preprocess_text(file_path, raw))
            slide_text = "\n".join(text_runs)
            if slide_text:
                docs.append(Document(
                    text=slide_text,
                    metadata={
                        "file_path": file_path,
                        "slide_number": slide_idx,
                        "file_type": "pptx",
                    }
                ))
        return docs

# ------------------------------
# CUSTOM PDF READER
# ------------------------------
class PDFTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        docs: list[Document] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                raw = page.extract_text()
                if raw:
                    cleaned = preprocess_text(file_path, raw)
                    docs.append(Document(
                        text=cleaned,
                        metadata={
                            "file_path": file_path,
                            "page_number": page.page_number, 
                            "file_type": "pdf",
                        }
                    ))
        return docs

# ------------------------------
# CUSTOM NOTEBOOK READER AND PARSER
# ------------------------------
class NotebookReader(BaseReader):
    def __init__(self, include_outputs=True, include_metadata=True):
        self.ipynb_reader = IPYNBReader()
        self.include_outputs = include_outputs
        self.include_metadata = include_metadata
    
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        try:
            # Use LlamaIndex's built-in IPYNB reader
            docs = self.ipynb_reader.load_data(file_path)
            
            # Enhanced processing for notebook content
            processed_docs = []
            for doc in docs:
                # Parse notebook content more intelligently
                notebook_content = self._parse_notebook_content(doc.text, file_path)
                
                # Update metadata with notebook-specific information
                enhanced_metadata = {
                    "file_path": file_path,
                    "file_name": os.path.basename(file_path),  # <-- Add this line
                    "file_type": "notebook",
                    "folder_name": os.path.basename(os.path.dirname(file_path)),
                    "notebook_cells": self._count_cells(doc.text),
                    "has_code": self._has_code_cells(doc.text),
                    "has_markdown": self._has_markdown_cells(doc.text)
                }
                
                # Merge with existing metadata
                if hasattr(doc, 'metadata') and doc.metadata:
                    enhanced_metadata.update(doc.metadata)
                
                # Create processed document
                processed_docs.append(Document(
                    text=notebook_content,
                    metadata=enhanced_metadata
                ))
            
            return processed_docs
            
        except Exception as e:
            print(f"❌ Error reading notebook {file_path}: {e}")
            return []
    
    def _parse_notebook_content(self, content: str, file_path: str) -> str:
        """Enhanced parsing of notebook content"""
        # Clean and preprocess the notebook content
        cleaned_content = preprocess_text(file_path, content)
        
        # Add notebook structure markers for better semantic understanding
        if "```python" in cleaned_content or "```" in cleaned_content:
            # This indicates code blocks are present
            cleaned_content = self._enhance_code_blocks(cleaned_content)
        
        return cleaned_content
    
    def _enhance_code_blocks(self, content: str) -> str:
        """Add semantic markers to code blocks"""
        # Add context markers for better chunking
        content = re.sub(r'```python\n', '\n--- CODE CELL START ---\n```python\n', content)
        content = re.sub(r'```\n', '\n```\n--- CODE CELL END ---\n', content)
        return content
    
    def _count_cells(self, content: str) -> int:
        """Count approximate number of cells"""
        return content.count('```') // 2 + content.count('# ')
    
    def _has_code_cells(self, content: str) -> bool:
        """Check if notebook has code cells"""
        return '```python' in content or '```' in content
    
    def _has_markdown_cells(self, content: str) -> bool:
        """Check if notebook has markdown cells"""
        return '#' in content or '*' in content or '[' in content

# ------------------------------
# SPECIALIZED NOTEBOOK PARSER
# ------------------------------
class NotebookAwareParser:
    """Custom parser that handles notebooks differently from other documents"""
    
    def __init__(self, embed_model):
        self.embed_model = embed_model
        
        # Create different parsers for different content types
        self.semantic_splitter = SemanticSplitterNodeParser(
            buffer_size=1,
            breakpoint_percentile_threshold=95,
            embed_model=embed_model
        )
        
        self.sentence_splitter = SentenceSplitter(
            chunk_size=512,
            chunk_overlap=50
        )
        
        # Code splitter for better handling of code content
        if CODE_SPLITTER_AVAILABLE:
            self.code_splitter = CodeSplitter(
                language="python",
                chunk_lines=40,
                chunk_lines_overlap=15,
                max_chars=1500
            )
        else:
            self.code_splitter = None
    
    def parse_documents(self, documents):
        """Parse documents using appropriate parser based on content type"""
        all_nodes = []
        
        for doc in documents:
            file_type = doc.metadata.get('file_type', 'unknown')
            
            if file_type == 'notebook':
                nodes = self._parse_notebook(doc)
            elif self._is_code_file(doc):
                nodes = self._parse_code_file(doc)
            else:
                nodes = self._parse_regular_document(doc)
            
            all_nodes.extend(nodes)
        
        return all_nodes
    
    def _parse_notebook(self, doc):
        """Special parsing for notebook documents"""
        print(f"📓 Parsing notebook: {doc.metadata.get('file_name', 'unknown')}")
        
        # Split notebook content by cell markers if present
        content = doc.text
        
        if "--- CODE CELL START ---" in content:
            # Parse code and markdown sections separately
            sections = re.split(r'--- CODE CELL (?:START|END) ---', content)
            nodes = []
            
            for i, section in enumerate(sections):
                if section.strip():
                    if '```python' in section and self.code_splitter:
                        # Use code splitter for code sections
                        temp_doc = Document(text=section, metadata=doc.metadata.copy())
                        code_nodes = self.code_splitter.get_nodes_from_documents([temp_doc])
                        
                        # Add cell information to metadata
                        for node in code_nodes:
                            node.metadata['cell_type'] = 'code'
                            node.metadata['section_index'] = i
                        
                        nodes.extend(code_nodes)
                    else:
                        # Use semantic splitter for markdown sections
                        temp_doc = Document(text=section, metadata=doc.metadata.copy())
                        semantic_nodes = self.semantic_splitter.get_nodes_from_documents([temp_doc])
                        
                        # Add cell information to metadata
                        for node in semantic_nodes:
                            node.metadata['cell_type'] = 'markdown'
                            node.metadata['section_index'] = i
                        
                        nodes.extend(semantic_nodes)
            
            return nodes
        else:
            # Fallback to semantic splitting
            return self.semantic_splitter.get_nodes_from_documents([doc])
    
    def _parse_code_file(self, doc):
        """Parse code files with code splitter if available"""
        if self.code_splitter:
            return self.code_splitter.get_nodes_from_documents([doc])
        else:
            return self.sentence_splitter.get_nodes_from_documents([doc])
    
    def _parse_regular_document(self, doc):
        """Parse regular documents with semantic splitter"""
        try:
            return self.semantic_splitter.get_nodes_from_documents([doc])
        except Exception as e:
            print(f"⚠️ Semantic parsing failed for {doc.metadata.get('file_name', 'unknown')}, using sentence splitter: {e}")
            return self.sentence_splitter.get_nodes_from_documents([doc])
    
    def _is_code_file(self, doc):
        """Check if document is a code file"""
        file_path = doc.metadata.get('file_path', '')
        if isinstance(file_path, Path):
            file_path = str(file_path)
        code_extensions = {'.py', '.js', '.java', '.cpp', '.c', '.cs', '.html', '.css', '.php', '.rb'}
        return any(file_path.lower().endswith(ext) for ext in code_extensions)

# ------------------------------
# MODEL & LLM SETTINGS
# ------------------------------
try:
    model_name = "sentence-transformers/all-MiniLM-L6-v2"
    Settings.embed_model = HuggingFaceEmbedding(model_name=model_name)
    print(f"✅ Model {model_name} loaded successfully.")
except Exception as e:
    print(f"❌ Error loading embedding model: {e}")
    exit()

Settings.llm = Ollama(model="llama3.1:latest", request_timeout=120.0)

# ------------------------------
# TEXT PREPROCESSING FUNCTION
# ------------------------------
def preprocess_text(file_path, text):
    code_extensions = {".py", ".java", ".cpp", ".js", ".c", ".cs", ".html", ".css", ".php", ".rb", ".ipynb"}
    text_extensions = {".pdf", ".docx", ".pptx", ".txt"}

    ext = os.path.splitext(file_path)[-1].lower()

    if ext in text_extensions or ext in code_extensions:
        text = clean_text(text)
    return text

def clean_text(text):
    text = unicodedata.normalize("NFKD", text)

    # Preserve email addresses
    email_pattern = r'([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})'
    emails = re.findall(email_pattern, text)
    for i, email in enumerate(emails):
        text = text.replace(email, f'EMAIL_PLACEHOLDER_{i}')

    # Preserve URLs
    url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
    urls = re.findall(url_pattern, text)
    for i, url in enumerate(urls):
        text = text.replace(url, f'URL_PLACEHOLDER_{i}')

    # Remove table of contents artifacts (long sequences of dots)
    text = re.sub(r'\.{5,}', ' ', text)  # Remove sequences of 5 or more dots

    # Remove excessive spaces and unnecessary formatting
    text = re.sub(r'\s+', ' ', text).strip()  

    # Restore emails and URLs
    for i, email in enumerate(emails):
        text = text.replace(f'EMAIL_PLACEHOLDER_{i}', email)
    for i, url in enumerate(urls):
        text = text.replace(f'URL_PLACEHOLDER_{i}', url)

    return text

# ------------------------------
# LOAD DOCUMENTS FROM MULTIPLE DIRECTORIES
# ------------------------------
def load_documents_from_directories(directories, file_extensions=None):
    """Load documents from multiple directories"""
    if file_extensions is None:
        file_extensions = ['.pptx', '.ipynb', '.docx', '.csv', '.jpeg', '.pdf', '.png', '.py']
    
    all_docs = []
    
    for directory in directories:
        if not os.path.exists(directory):
            print(f"⚠️ Directory {directory} does not exist, skipping...")
            continue
            
        print(f"📁 Loading documents from: {directory}")
        
        try:
            reader = SimpleDirectoryReader(
                input_dir=directory,
                required_exts=file_extensions,
                file_extractor={
                    ".pptx": PPTXTextOnlyReader(), 
                    ".pdf": PDFTextOnlyReader(),
                    ".ipynb": NotebookReader()
                }, 
                recursive=True
            )
            docs = reader.load_data()
            
            if docs:
                print(f"✅ Loaded {len(docs)} documents from {directory}")
                all_docs.extend(docs)
            else:
                print(f"⚠️ No documents found in {directory}")
                
        except Exception as e:
            print(f"❌ Error loading documents from {directory}: {e}")
    
    return all_docs

# Define your directories - modify these paths as needed
doc_directories = [
    "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules/",
    # Add more directories here as needed
    # "/path/to/notebooks/directory1/",
    # "/path/to/notebooks/directory2/",
]

# Load documents from all directories
docs = load_documents_from_directories(doc_directories)

if not docs:
    print("❌ No documents found! Check the paths and file extensions.")
    exit()

print(f"✅ Total loaded documents: {len(docs)}")

# ------------------------------
# ENHANCED SEMANTIC CHUNKING SETUP
# ------------------------------
def create_notebook_aware_parser():
    """Create a notebook-aware parser that handles different content types appropriately"""
    try:
        notebook_parser = NotebookAwareParser(embed_model=Settings.embed_model)
        print("✅ Notebook-aware parser created successfully")
        return notebook_parser
    except Exception as e:
        print(f"⚠️ Error creating notebook-aware parser: {e}")
        # Fallback to basic semantic splitter
        return SemanticSplitterNodeParser(
            buffer_size=1,
            breakpoint_percentile_threshold=95,
            embed_model=Settings.embed_model
        )

# Create the enhanced parser
notebook_parser = create_notebook_aware_parser()

# ------------------------------
# PROCESS DOCUMENTS WITH NOTEBOOK-AWARE PARSING
# ------------------------------
def process_documents_with_notebook_aware_parsing(documents, parser):
    """Process documents using notebook-aware parsing"""
    processed_docs = []
    
    print("🔄 Processing documents with notebook-aware parsing...")
    
    # Group documents by type for better processing
    notebooks = [doc for doc in documents if doc.metadata.get('file_type') == 'notebook']
    other_docs = [doc for doc in documents if doc.metadata.get('file_type') != 'notebook']
    
    print(f"📓 Found {len(notebooks)} notebook documents")
    print(f"📄 Found {len(other_docs)} other documents")
    
    # Process all documents
    if hasattr(parser, 'parse_documents'):
        # Use the custom notebook-aware parser
        nodes = parser.parse_documents(documents)
        
        # Convert nodes to Document objects with enhanced metadata
        for i, node in enumerate(nodes):
            # Enhance metadata
            enhanced_metadata = node.metadata.copy()
            enhanced_metadata.update({
                "chunk_index": i,
                "num_tokens": len(node.text.split()),
                "num_chars": len(node.text),
                "parsing_method": "notebook_aware"
            })
            
            # Add notebook-specific metadata if available
            if 'cell_type' in node.metadata:
                enhanced_metadata['cell_type'] = node.metadata['cell_type']
            if 'section_index' in node.metadata:
                enhanced_metadata['section_index'] = node.metadata['section_index']
            
            processed_docs.append(Document(
                text=node.text,
                metadata=enhanced_metadata
            ))
    else:
        # Fallback to standard processing
        print("⚠️ Using fallback processing method")
        for doc in documents:
            try:
                # Apply standard chunking
                chunks = parser.get_nodes_from_documents([doc])
                
                for i, chunk in enumerate(chunks):
                    enhanced_metadata = doc.metadata.copy()
                    enhanced_metadata.update({
                        "chunk_index": i,
                        "num_tokens": len(chunk.text.split()),
                        "num_chars": len(chunk.text),
                        "parsing_method": "standard"
                    })
                    
                    processed_docs.append(Document(
                        text=chunk.text,
                        metadata=enhanced_metadata
                    ))
                    
            except Exception as e:
                print(f"⚠️ Error processing document {doc.metadata.get('file_path', 'unknown')}: {e}")
                continue
    
    return processed_docs

# Process documents with notebook-aware parsing
print("🔄 Processing documents with enhanced notebook-aware parsing...")
document_objects = process_documents_with_notebook_aware_parsing(docs, notebook_parser)
print(f"✅ Created {len(document_objects)} intelligently parsed chunks")

# ------------------------------
# CHUNK VISUALIZATION AND STORAGE
# ------------------------------
def print_and_save_chunks(chunks, output_file="chunks_output.txt"):
    """Print chunks to console and save to text file"""
    try:
        print("\n" + "="*80)
        print("📄 CHUNK ANALYSIS AND PREVIEW")
        print("="*80)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            # Write header
            f.write("SEMANTIC CHUNKS ANALYSIS\n")
            f.write("="*80 + "\n")
            f.write(f"Total Chunks: {len(chunks)}\n")
            f.write(f"Generated on: {os.popen('date').read().strip()}\n")
            f.write("="*80 + "\n\n")
            
            # Group chunks by file for better organization
            chunks_by_file = {}
            for i, chunk in enumerate(chunks):
                file_name = chunk.metadata.get('file_name', 'unknown')
                if file_name not in chunks_by_file:
                    chunks_by_file[file_name] = []
                chunks_by_file[file_name].append((i, chunk))
            
            # Print summary
            print(f"📊 CHUNK SUMMARY:")
            print(f"   Total chunks: {len(chunks)}")
            print(f"   Files processed: {len(chunks_by_file)}")
            print(f"   Output file: {output_file}")
            
            f.write(f"CHUNK SUMMARY:\n")
            f.write(f"Total chunks: {len(chunks)}\n")
            f.write(f"Files processed: {len(chunks_by_file)}\n\n")
            
            # Analyze chunk statistics
            chunk_sizes = [chunk.metadata.get('num_chars', 0) for chunk in chunks]
            token_counts = [chunk.metadata.get('num_tokens', 0) for chunk in chunks]
            
            avg_chars = sum(chunk_sizes) / len(chunk_sizes) if chunk_sizes else 0
            avg_tokens = sum(token_counts) / len(token_counts) if token_counts else 0
            
            print(f"   Average chunk size: {avg_chars:.0f} characters, {avg_tokens:.0f} tokens")
            f.write(f"Average chunk size: {avg_chars:.0f} characters, {avg_tokens:.0f} tokens\n")
            
            # File breakdown
            print(f"\n📁 FILES BREAKDOWN:")
            f.write(f"\nFILES BREAKDOWN:\n")
            for file_name, file_chunks in chunks_by_file.items():
                print(f"   {file_name}: {len(file_chunks)} chunks")
                f.write(f"{file_name}: {len(file_chunks)} chunks\n")
            
            f.write("\n" + "="*80 + "\n\n")
            
            # Process each file
            for file_name, file_chunks in chunks_by_file.items():
                print(f"\n📄 FILE: {file_name}")
                print("-" * 60)
                
                f.write(f"FILE: {file_name}\n")
                f.write("-" * 60 + "\n")
                
                for chunk_idx, (global_idx, chunk) in enumerate(file_chunks):
                    # Print to console (first few chunks of each file)
                    if chunk_idx < 3:  # Show first 3 chunks of each file in console
                        print(f"\n🔸 Chunk {chunk_idx + 1}/{len(file_chunks)} (Global ID: {global_idx})")
                        print(f"   📍 Source: {chunk.metadata.get('folder_name', 'N/A')}")
                        
                        # Show specific location info
                        location_info = []
                        if 'page_number' in chunk.metadata:
                            location_info.append(f"Page {chunk.metadata['page_number']}")
                        if 'slide_number' in chunk.metadata:
                            location_info.append(f"Slide {chunk.metadata['slide_number']}")
                        if 'cell_number' in chunk.metadata:
                            location_info.append(f"Cell {chunk.metadata['cell_number']}")
                        
                        if location_info:
                            print(f"   📍 Location: {', '.join(location_info)}")
                        
                        print(f"   📏 Size: {chunk.metadata.get('num_chars', 0)} chars, {chunk.metadata.get('num_tokens', 0)} tokens")
                        print(f"   🔤 Content preview:")
                        
                        # Show preview of content
                        content_preview = chunk.text[:200] + "..." if len(chunk.text) > 200 else chunk.text
                        print(f"      {content_preview}")
                        
                    elif chunk_idx == 3 and len(file_chunks) > 3:
                        print(f"\n   ... and {len(file_chunks) - 3} more chunks (see {output_file} for full details)")
                    
                    # Write all chunks to file
                    f.write(f"\nCHUNK {chunk_idx + 1}/{len(file_chunks)} (Global ID: {global_idx})\n")
                    f.write(f"Source Folder: {chunk.metadata.get('folder_name', 'N/A')}\n")
                    f.write(f"File Type: {chunk.metadata.get('file_type', 'unknown')}\n")
                    
                    # Location information
                    if 'page_number' in chunk.metadata:
                        f.write(f"Page Number: {chunk.metadata['page_number']}\n")
                    if 'slide_number' in chunk.metadata:
                        f.write(f"Slide Number: {chunk.metadata['slide_number']}\n")
                    if 'cell_number' in chunk.metadata:
                        f.write(f"Cell Number: {chunk.metadata['cell_number']}\n")
                    
                    f.write(f"Chunk Index: {chunk.metadata.get('chunk_index', 'N/A')}\n")
                    f.write(f"Characters: {chunk.metadata.get('num_chars', 0)}\n")
                    f.write(f"Tokens: {chunk.metadata.get('num_tokens', 0)}\n")
                    f.write(f"Content:\n{'-' * 40}\n")
                    f.write(f"{chunk.text}\n")
                    f.write(f"{'-' * 40}\n\n")
                
                f.write(f"\n{'=' * 80}\n\n")
        
        print(f"\n✅ Chunk analysis completed!")
        print(f"📁 Full details saved to: {output_file}")
        print(f"📊 Summary: {len(chunks)} chunks from {len(chunks_by_file)} files")
        
    except Exception as e:
        print(f"❌ Error saving chunks: {e}")

# Print and save chunks
print_and_save_chunks(document_objects)

# ------------------------------
# DOCUMENT PROCESSING PIPELINE
# ------------------------------
pipeline = IngestionPipeline(
    transformations=[
        Settings.embed_model  # Only apply the embedding model to pre-chunked text
    ],
)

try:
    # Process the semantically chunked documents
    nodes = pipeline.run(documents=document_objects)
    
    if not nodes:
        print("❌ No nodes were created. Check document parsing.")
        exit()
    print(f"✅ {len(nodes)} document nodes created and ready for ChromaDB.")

    # Add nodes to ChromaDB
    chroma_path = "./chroma_db"
    chroma_client = PersistentClient(path=chroma_path)
    collection = chroma_client.get_or_create_collection("document_chunks")

    vector_store = ChromaVectorStore(chroma_client, collection_name="document_chunks")
    
    for idx, node in enumerate(nodes):
        # Sanitize metadata values: convert any Path → str
        clean_meta = {
            key: str(val) if isinstance(val, Path) else val
            for key, val in node.metadata.items()
        }

        # Add to Chroma, using the cleaned metadata
        collection.add(
            ids=[str(idx)],
            documents=[node.text],
            metadatas=[clean_meta]
        )
    
    print(f"✅ {len(nodes)} nodes stored in ChromaDB with semantic chunking.")

except Exception as e:
    print(f"❌ Error during ingestion pipeline: {e}")
    exit()

# ------------------------------
# CREATE VECTOR STORE INDEX
# ------------------------------
# Clean metadata for all nodes
for node in nodes:
    for k, v in list(node.metadata.items()):
        if isinstance(v, Path):
            node.metadata[k] = str(v)

try:
    index = VectorStoreIndex(nodes, vector_store=vector_store)
    print("✅ Vector store index created successfully with semantic chunking.")
    
    persist_dir = "./persisted_index"
    os.makedirs(persist_dir, exist_ok=True)
    index.storage_context.persist(persist_dir=persist_dir)
    print(f"✅ Index persisted to {persist_dir}")
except Exception as e:
    print(f"❌ Error creating VectorStoreIndex: {e}")
    exit()

# ------------------------------
# CREATE CHAT ENGINE & PROCESS QUERY
# ------------------------------
class RAGQueryEngine(CustomQueryEngine):
    """RAG Query Engine for custom retrieval and response synthesis."""

    retriever: BaseRetriever
    response_synthesizer: BaseSynthesizer

    def custom_query(self, query_str: str):
        nodes = self.retriever.retrieve(query_str)
        response_obj = self.response_synthesizer.synthesize(query_str, nodes)
        return response_obj

retriever = index.as_retriever(similarity_top_k=5)  # Retrieve more chunks for better context
synthesizer = get_response_synthesizer(response_mode="compact")

query_engine = RAGQueryEngine(
    retriever=retriever, response_synthesizer=synthesizer
)

# Test the system
print("\n" + "="*50)
print("🚀 Testing the enhanced RAG system...")
print("="*50)

response = query_engine.query("Write a code to find a factorial for a number?")
print("Query: Write a code to find a factorial for a number?")
print(f"Response: {response}")

# Additional test for notebook content
print("\n" + "-"*50)
notebook_response = query_engine.query("Show me any Python functions or code examples from notebooks")
print("Query: Show me any Python functions or code examples from notebooks")
print(f"Response: {notebook_response}")

print("\n✅ Enhanced RAG system with semantic chunking and notebook support is ready!")

# ------------------------------
# ADDITIONAL CHUNK ANALYSIS FUNCTIONS
# ------------------------------
def analyze_chunk_distribution(chunks, output_file="chunk_analysis.txt"):
    """Analyze and save detailed chunk distribution statistics"""
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("DETAILED CHUNK ANALYSIS\n")
            f.write("="*80 + "\n\n")
            
            # File type distribution
            file_types = {}
            folder_distribution = {}
            size_distribution = {'small': 0, 'medium': 0, 'large': 0}
            
            for chunk in chunks:
                # File type analysis
                file_type = chunk.metadata.get('file_type', 'unknown')
                file_types[file_type] = file_types.get(file_type, 0) + 1
                
                # Folder analysis
                folder = chunk.metadata.get('folder_name', 'unknown')
                folder_distribution[folder] = folder_distribution.get(folder, 0) + 1
                
                # Size analysis
                chars = chunk.metadata.get('num_chars', 0)
                if chars < 200:
                    size_distribution['small'] += 1
                elif chars < 800:
                    size_distribution['medium'] += 1
                else:
                    size_distribution['large'] += 1
            
            # Write analysis
            f.write("FILE TYPE DISTRIBUTION:\n")
            for file_type, count in sorted(file_types.items()):
                percentage = (count / len(chunks)) * 100
                f.write(f"  {file_type}: {count} chunks ({percentage:.1f}%)\n")
            
            f.write(f"\nFOLDER DISTRIBUTION:\n")
            for folder, count in sorted(folder_distribution.items()):
                percentage = (count / len(chunks)) * 100
                f.write(f"  {folder}: {count} chunks ({percentage:.1f}%)\n")
            
            f.write(f"\nCHUNK SIZE DISTRIBUTION:\n")
            for size_type, count in size_distribution.items():
                percentage = (count / len(chunks)) * 100
                f.write(f"  {size_type} (<200 chars, 200-800 chars, >800 chars): {count} chunks ({percentage:.1f}%)\n")
            
        print(f"📊 Detailed analysis saved to: {output_file}")
        
    except Exception as e:
        print(f"❌ Error creating chunk analysis: {e}")

def save_chunks_by_type(chunks, output_dir="chunks_by_type"):
    """Save chunks grouped by file type in separate files"""
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        chunks_by_type = {}
        for chunk in chunks:
            file_type = chunk.metadata.get('file_type', 'unknown')
            if file_type not in chunks_by_type:
                chunks_by_type[file_type] = []
            chunks_by_type[file_type].append(chunk)
        
        for file_type, type_chunks in chunks_by_type.items():
            output_file = os.path.join(output_dir, f"{file_type}_chunks.txt")
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"CHUNKS FROM {file_type.upper()} FILES\n")
                f.write("="*60 + "\n")
                f.write(f"Total chunks: {len(type_chunks)}\n\n")
                
                for i, chunk in enumerate(type_chunks):
                    f.write(f"CHUNK {i+1}\n")
                    f.write(f"File: {chunk.metadata.get('file_name', 'unknown')}\n")
                    f.write(f"Size: {chunk.metadata.get('num_chars', 0)} chars\n")
                    f.write(f"Content:\n{'-'*40}\n")
                    f.write(f"{chunk.text}\n")
                    f.write(f"{'-'*40}\n\n")
        
        print(f"📁 Chunks by type saved to: {output_dir}/")
        
    except Exception as e:
        print(f"❌ Error saving chunks by type: {e}")

# Run additional analysis
print("\n🔍 Running additional chunk analysis...")
analyze_chunk_distribution(document_objects)
save_chunks_by_type(document_objects)

print(f"\n🎉 All chunk analysis completed!")
print(f"📄 Main chunks file: chunks_output.txt")
print(f"📊 Analysis file: chunk_analysis.txt") 
print(f"📁 Type-separated chunks: chunks_by_type/ directory")